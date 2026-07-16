from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
bullet_slide_layout = prs.slide_layouts[1]

def add_title_slide(title_text, subtitle_text):
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = title_text
    subtitle.text = subtitle_text

def add_bullet_slide(title_text, bullet_points):
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = title_text
    
    tf = body_shape.text_frame
    
    first = True
    for point in bullet_points:
        if isinstance(point, tuple):
            text, level = point
        else:
            text, level = point, 0
            
        if first:
            tf.text = text
            first = False
        else:
            p = tf.add_paragraph()
            p.text = text
            p.level = level

# Slide 1
add_title_slide("Go Chicken", "Pitch Deck Storyboard & Context Report")

# Slide 2
add_bullet_slide(
    "The Reality of Wholesale Poultry Today",
    [
        "Managing wholesale poultry relies heavily on manual, unstructured communication",
        ("Creates massive operational bottlenecks, pricing disputes, and lost revenue.", 1),
        "Latency",
        ("Takes up to 30 minutes to negotiate, verify, and confirm.", 1),
        "Errors",
        ("Misread WhatsApps lead to duplicate orders and misrouted trucks.", 1),
        "Pricing Disputes",
        ("Volatile live-bird prices lead to margin loss when quoting outdated prices.", 1)
    ]
)

# Slide 3
add_bullet_slide(
    "Meet Raj, The Poultry Wholesaler",
    [
        "Raj, The Poultry Wholesaler (Main Boss)",
        ("Goal: Maximize profitability and scale operations.", 1),
        ("Reality: Stressed, blind to real-time inventory, unaware of exact Khata debt.", 1),
        "Other Stakeholders Impacted:",
        ("Retailer: Wasting time negotiating, unaware of their exact debt.", 1),
        ("Delivery Driver: Verbal dispatches leading to miscommunications.", 1)
    ]
)

# Slide 4
add_bullet_slide(
    "What if WhatsApp could run the business?",
    [
        "Vision: Bring enterprise software to where the business already happens",
        ("No need for retailers to download a clunky, generic ERP app.", 1),
        "The Flow:",
        ("Retailer sends a text → AI understands it", 1),
        ("Backend prices it → Dashboard updates live", 1),
        ("Ledger is reconciled", 1),
        "All under 1 second. Zero friction."
    ]
)

# Slide 5
add_bullet_slide(
    "The Live Demo (The \"Wow\" Moment)",
    [
        "1. The Retailer Phone Screen",
        ("Retailer texts: \"Need 50kg chicken, what's today's rate?\"", 1),
        ("Bot replies instantly (< 1s) with Quote, Khata balance & buttons.", 1),
        "2. The Wholesaler Dashboard",
        ("Monochromatic Next.js Dashboard updates instantly via SSE (no refresh).", 1),
        ("Order appears, inventory drops, analytics adjust.", 1),
        "3. The Pricing Magic",
        ("Raj updates price on dashboard.", 1),
        ("Retailer asks for rate again -> bot quotes newly updated price.", 1)
    ]
)

# Slide 6
add_bullet_slide(
    "Technical Architecture",
    [
        "Highly scalable, event-driven enterprise engine",
        "The Data Pipeline",
        ("Meta API → FastAPI → AI (Groq) → Event Backend → Supabase → SSE → Next.js", 1),
        "Frontend: Next.js",
        ("B&W Enterprise UI, SSE live updates, Optimistic Rollbacks.", 1),
        "Backend: FastAPI (Python)",
        ("RESTful API and direct WhatsApp Webhook handler for sub-second latency.", 1),
        "Database & AI",
        ("PostgreSQL with pgvector on Supabase.", 1),
        ("Groq (production) / Ollama (local) with 100% Regex Fallback.", 1)
    ]
)

# Slide 7
add_bullet_slide(
    "Enterprise Engineering Highlights",
    [
        "Advanced software engineering patterns",
        ("Hierarchical Pricing Engine: Dynamic pricing & Immutable Quote Snapshots.", 1),
        ("Quote-to-Order Conversion: Deterministic flow (no AI hallucinations).", 1),
        ("CQRS & Analytics: Projection Rebuilder for lightning-fast metrics.", 1),
        ("Event-Driven Architecture: Transactional Outbox pattern.", 1),
        ("Immutable Khata Ledger: Append-only events for perfect auditability.", 1),
        ("Server-Sent Events (SSE): Truly reactive SPA experience.", 1),
        ("Multi-Tenant Design: Strict server-side isolation.", 1)
    ]
)

# Slide 8
add_bullet_slide(
    "Business Value",
    [
        "Zero-Latency Ordering",
        ("Order processing time drops from 30 minutes to 1 second.", 1),
        "Margin Protection",
        ("Dynamic Pricing API & Immutable Quotes prevent unprofitable sales.", 1),
        "Zero Khata Errors",
        ("Automated ledger updates prevent disputes and bad debt.", 1),
        "Seamless Adoption",
        ("Retailers don't download an app; they just use WhatsApp.", 1)
    ]
)

# Slide 9
add_bullet_slide(
    "The Vision (Roadmap)",
    [
        "Immediately after Hackathon",
        ("Deploy AI Forecaster Engine using pgvector to predict demand spikes.", 1),
        "6 Months",
        ("Integrate Fleet Route Optimization (Google Maps API) for delivery trucks.", 1),
        "1 Year",
        ("Expand into a B2B marketplace to trade surplus inventory directly.", 1)
    ]
)

output_path = r"d:\Go Chicken\Go_Chicken_Pitch_Deck.pptx"
prs.save(output_path)
print(f"Presentation saved to {output_path}")
